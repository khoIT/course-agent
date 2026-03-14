#!/usr/bin/env python3
"""Convert slide markdown files to .pptx presentations.

Usage: python convert-slides-to-pptx.py input.md output.pptx

Markdown format expected:
- First line: # Session Title
- Slides separated by ---
- Slide titles: ## Slide N: Title
- Speaker notes after **Ghi chu giang vien:**
- Bullet points as - items
- Tables, code blocks, and blockquotes preserved as text
"""

import sys
import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# --- Constants ---
SLIDE_WIDTH = Emu(12192000)   # 16:9 widescreen
SLIDE_HEIGHT = Emu(6858000)
TITLE_FONT_SIZE = Pt(28)
BODY_FONT_SIZE = Pt(16)
SMALL_FONT_SIZE = Pt(14)
TITLE_COLOR = RGBColor(0x1A, 0x1A, 0x2E)  # Dark navy
BODY_COLOR = RGBColor(0x33, 0x33, 0x33)
ACCENT_COLOR = RGBColor(0x2D, 0x5F, 0x9A)  # Blue accent
BG_COLOR = RGBColor(0xFF, 0xFF, 0xFF)


def parse_markdown(md_text: str):
    """Parse markdown into session title and list of slide dicts."""
    lines = md_text.strip().split('\n')

    # Extract session title from first line
    session_title = ''
    content_start = 0
    for i, line in enumerate(lines):
        if line.startswith('# '):
            session_title = line[2:].strip()
            content_start = i + 1
            break

    # Split remaining content on --- separators
    remaining = '\n'.join(lines[content_start:])
    raw_slides = re.split(r'\n---\n', remaining)

    slides = []
    for raw in raw_slides:
        raw = raw.strip()
        if not raw:
            continue
        slide = parse_slide(raw)
        if slide:
            slides.append(slide)

    return session_title, slides


def parse_slide(raw: str):
    """Parse a single slide's markdown into structured data."""
    lines = raw.split('\n')

    # Extract title from ## Slide N: Title
    title = ''
    body_lines = []
    notes = ''

    title_found = False
    in_notes = False
    notes_lines = []

    for line in lines:
        # Check for slide title
        m = re.match(r'^##\s+Slide\s+\d+\s*:\s*(.*)', line)
        if m and not title_found:
            title = m.group(1).strip()
            # Clean up title: remove [...] brackets that are layout instructions
            title = re.sub(r'\s*\[.*?\]\s*$', '', title).strip()
            # Remove leading/trailing --
            title = re.sub(r'^--\s*', '', title).strip()
            title = re.sub(r'\s*--\s*$', '', title).strip()
            title_found = True
            continue

        # Check for speaker notes
        if re.match(r'^\*\*Ghi ch', line):
            in_notes = True
            # Extract text after the marker
            note_text = re.sub(r'^\*\*Ghi ch[^*]*\*\*\s*', '', line)
            notes_lines.append(note_text)
            continue

        if in_notes:
            # Notes continue until a blank line or new section
            if line.strip() == '' or line.startswith('##') or line.startswith('---'):
                in_notes = False
            else:
                notes_lines.append(line)
            continue

        body_lines.append(line)

    notes = ' '.join(notes_lines).strip()
    body = '\n'.join(body_lines).strip()

    if not title and not body:
        return None

    return {
        'title': title if title else '(Slide)',
        'body': body,
        'notes': notes,
    }


def clean_markdown_formatting(text: str) -> str:
    """Remove markdown formatting for plain text display."""
    # Remove bold markers
    text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
    # Remove italic markers
    text = re.sub(r'\*(.*?)\*', r'\1', text)
    # Remove inline code
    text = re.sub(r'`(.*?)`', r'\1', text)
    # Remove image/link syntax
    text = re.sub(r'\[([^\]]*)\]\([^)]*\)', r'\1', text)
    return text


def add_title_slide(prs, session_title: str):
    """Add the title slide for the session."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = TITLE_COLOR

    # Title text box - centered
    left = Inches(1)
    top = Inches(2)
    width = Inches(11.33)
    height = Inches(2.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = session_title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER


def process_body_content(tf, body: str):
    """Process slide body content into text frame paragraphs."""
    lines = body.split('\n')
    first_para = True

    i = 0
    in_code_block = False
    in_table = False
    table_lines = []

    while i < len(lines):
        line = lines[i]

        # Handle code blocks
        if line.strip().startswith('```'):
            in_code_block = not in_code_block
            if not in_code_block and not first_para:
                # End of code block
                pass
            i += 1
            continue

        if in_code_block:
            p = tf.paragraphs[0] if first_para else tf.add_paragraph()
            first_para = False
            p.text = line
            p.font.size = SMALL_FONT_SIZE
            p.font.color.rgb = BODY_COLOR
            p.font.name = 'Consolas'
            i += 1
            continue

        # Handle table rows
        if '|' in line and line.strip().startswith('|'):
            # Skip separator rows
            if re.match(r'^[\s|:-]+$', line):
                i += 1
                continue
            # Parse table cells
            cells = [c.strip() for c in line.split('|')[1:-1]]
            cell_text = '  |  '.join(cells)
            cell_text = clean_markdown_formatting(cell_text)
            p = tf.paragraphs[0] if first_para else tf.add_paragraph()
            first_para = False
            p.text = cell_text
            p.font.size = SMALL_FONT_SIZE
            p.font.color.rgb = BODY_COLOR
            i += 1
            continue

        # Skip empty lines
        if not line.strip():
            i += 1
            continue

        # Skip image placeholders [Hinh minh hoa: ...]
        if re.match(r'^\[.*?\]$', line.strip()):
            i += 1
            continue

        # Blockquotes
        if line.strip().startswith('>'):
            text = re.sub(r'^>\s*', '', line.strip())
            text = clean_markdown_formatting(text)
            p = tf.paragraphs[0] if first_para else tf.add_paragraph()
            first_para = False
            p.text = f'"{text}"'
            p.font.size = BODY_FONT_SIZE
            p.font.italic = True
            p.font.color.rgb = ACCENT_COLOR
            i += 1
            continue

        # Bold section headers (like **Demo truc tiep:** or **Buoc 1:**)
        m = re.match(r'^\*\*(.*?)\*\*\s*$', line.strip())
        if m:
            p = tf.paragraphs[0] if first_para else tf.add_paragraph()
            first_para = False
            p.text = m.group(1)
            p.font.size = BODY_FONT_SIZE
            p.font.bold = True
            p.font.color.rgb = ACCENT_COLOR
            p.space_before = Pt(8)
            i += 1
            continue

        # Bold labels with content (like **Label:** content)
        m = re.match(r'^\*\*(.*?)\*\*\s*(.*)', line.strip())
        if m and not line.strip().startswith('- '):
            label = m.group(1)
            content = clean_markdown_formatting(m.group(2))
            p = tf.paragraphs[0] if first_para else tf.add_paragraph()
            first_para = False
            run1 = p.runs[0] if p.runs else p.add_run()
            run1.text = label + ' '
            run1.font.size = BODY_FONT_SIZE
            run1.font.bold = True
            run1.font.color.rgb = ACCENT_COLOR
            if content:
                run2 = p.add_run()
                run2.text = content
                run2.font.size = BODY_FONT_SIZE
                run2.font.color.rgb = BODY_COLOR
            p.space_before = Pt(6)
            i += 1
            continue

        # Bullet points (- item or numbered 1. item)
        bm = re.match(r'^(\s*)[-*]\s+(.*)', line)
        nm = re.match(r'^(\s*)\d+\.\s+(.*)', line)
        if bm or nm:
            match = bm or nm
            indent = len(match.group(1))
            text = clean_markdown_formatting(match.group(2))
            p = tf.paragraphs[0] if first_para else tf.add_paragraph()
            first_para = False
            p.text = text
            p.font.size = BODY_FONT_SIZE
            p.font.color.rgb = BODY_COLOR
            p.level = min(indent // 2, 2)
            p.space_before = Pt(3)
            # Add bullet character
            if bm:
                p.text = '\u2022  ' + text
            else:
                p.text = text
            i += 1
            continue

        # Regular text
        text = clean_markdown_formatting(line.strip())
        if text:
            p = tf.paragraphs[0] if first_para else tf.add_paragraph()
            first_para = False
            p.text = text
            p.font.size = BODY_FONT_SIZE
            p.font.color.rgb = BODY_COLOR

        i += 1


def add_content_slide(prs, slide_data: dict):
    """Add a content slide with title, body, and notes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # White background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

    # Title bar - accent color strip at top
    from pptx.util import Emu as EmuUtil
    left = Inches(0)
    top = Inches(0)
    width = Inches(13.33)
    height = Inches(0.08)
    shape = slide.shapes.add_shape(
        1,  # rectangle
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_COLOR
    shape.line.fill.background()

    # Title text box
    title = slide_data['title']
    left = Inches(0.6)
    top = Inches(0.3)
    width = Inches(12)
    height = Inches(0.9)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = clean_markdown_formatting(title)
    p.font.size = TITLE_FONT_SIZE
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR

    # Body content
    if slide_data['body']:
        left = Inches(0.6)
        top = Inches(1.3)
        width = Inches(12)
        height = Inches(5.2)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = None

        process_body_content(tf, slide_data['body'])

    # Speaker notes
    if slide_data['notes']:
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = slide_data['notes']


def convert(input_path: str, output_path: str):
    """Convert a markdown slide file to .pptx."""
    md_text = Path(input_path).read_text(encoding='utf-8')
    session_title, slides = parse_markdown(md_text)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Title slide
    add_title_slide(prs, session_title)

    # Content slides
    for slide_data in slides:
        add_content_slide(prs, slide_data)

    prs.save(output_path)
    print(f"Created: {output_path} ({len(slides)} content slides + 1 title slide)")


if __name__ == '__main__':
    if len(sys.argv) != 3:
        print(f"Usage: {sys.argv[0]} input.md output.pptx")
        sys.exit(1)
    convert(sys.argv[1], sys.argv[2])
